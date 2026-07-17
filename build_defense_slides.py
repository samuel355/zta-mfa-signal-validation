#!/usr/bin/env python3
"""
Builds the MPhil thesis defense slide deck from the same corrected, real
results used in the thesis and manuscript (see
updated/reference_material/master_correction_map.md for the source of every
number used here). No fabricated figures.
"""
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
FIGURES_DIR = os.path.join(BASE_DIR, "updated", "figures")
OUT_PATH = os.path.join(BASE_DIR, "updated", "Defense_Slides.pptx")
METRICS = json.load(open(os.path.join(BASE_DIR, "scripts", "chapter4_metrics.json")))

# ── palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x14, 0x2850 >> 8 & 0xFF, 0x50)  # placeholder, overwritten below
NAVY   = RGBColor(0x14, 0x28, 0x50)
BLUE   = RGBColor(0x2A, 0x78, 0xD6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT  = RGBColor(0xF2, 0xF4, 0xF7)
GREEN  = RGBColor(0x1B, 0xAF, 0x7A)
RED    = RGBColor(0xE3, 0x49, 0x48)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                 color=GREY, align=PP_ALIGN.LEFT, italic=False, font="Calibri",
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_bar(slide, color=BLUE, top=Inches(1.15), height=Inches(0.06)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(2.2), height)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def slide_number(slide, n):
    add_textbox(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.5), Inches(0.7), Inches(0.35),
                str(n), size=11, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)


def title_slide(title, subtitle_lines, footer=""):
    s = add_slide()
    bg(s, NAVY)
    add_textbox(s, Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.8),
                title, size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    y = Inches(3.9)
    for line, size, bold, col in subtitle_lines:
        add_textbox(s, Inches(1.0), y, Inches(11.3), Inches(0.5), line,
                    size=size, bold=bold, color=col, align=PP_ALIGN.LEFT)
        y += Inches(0.45)
    if footer:
        add_textbox(s, Inches(1.0), SLIDE_H - Inches(0.9), Inches(11.3), Inches(0.4),
                    footer, size=13, color=RGBColor(0x9A, 0xB4, 0xD6), align=PP_ALIGN.LEFT)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.05), Inches(2.5), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    return s


def content_slide(title, n, eyebrow=None):
    s = add_slide()
    bg(s, WHITE)
    if eyebrow:
        add_textbox(s, Inches(0.6), Inches(0.35), Inches(8), Inches(0.35), eyebrow,
                    size=13, bold=True, color=BLUE)
    add_textbox(s, Inches(0.6), Inches(0.65), Inches(12.0), Inches(0.7), title,
                size=26, bold=True, color=NAVY)
    add_bar(s)
    slide_number(s, n)
    return s


def bullets(slide, items, left=Inches(0.7), top=Inches(1.5), width=Inches(11.9),
            height=Inches(5.4), size=17, color=GREY, gap=10, bold_lead=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = level
        prefix = "•  " if level == 0 else "-  "
        run = p.add_run()
        run.text = prefix + text
        run.font.size = Pt(size - level * 1.5)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_pic(slide, path, left, top, width=None, height=None):
    if not os.path.isfile(path):
        add_textbox(slide, left, top, width or Inches(4), height or Inches(3),
                    f"[missing figure: {os.path.basename(path)}]", size=12, italic=True, color=RED)
        return
    if width:
        slide.shapes.add_picture(path, left, top, width=width)
    else:
        slide.shapes.add_picture(path, left, top, height=height)


def stat_tile(slide, left, top, width, height, value, label, color=BLUE):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = color; box.line.width = Pt(1.25)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(8); tf.margin_bottom = Pt(8)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = value
    r1.font.size = Pt(30); r1.font.bold = True; r1.font.color.rgb = color; r1.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(13); r2.font.color.rgb = GREY; r2.font.name = "Calibri"


def add_table(slide, left, top, width, height, headers, rows, col_widths=None,
              header_color=NAVY, font_size=13):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gtable.columns[i].width = w
    for ci, h in enumerate(headers):
        cell = gtable.cell(0, ci)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(font_size); r.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtable.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(font_size); r.font.color.rgb = GREY
    return gtable


# ═══════════════════════════════════════════════════════════════════════════
n = 1

# 1. TITLE
title_slide(
    "A Multi-Source Context Validation Framework for\nAdaptive Zero Trust "
    "Multi-Factor Authentication\nin Remote Work Environments",
    [
        ("Samuel Osei Adu", 20, True, WHITE),
        ("Supervisor: Dr. Kornyo Oliver", 16, False, RGBColor(0xC7, 0xD6, 0xEC)),
        ("Department of Physical and Computational Science", 14, False, RGBColor(0xC7, 0xD6, 0xEC)),
        ("Kwame Nkrumah University of Science and Technology, Kumasi, Ghana", 14, False, RGBColor(0xC7, 0xD6, 0xEC)),
    ],
    footer="MPhil Thesis Defense — Cyber Security and Digital Forensics"
)
notes(prs.slides[-1],
    "Good [morning/afternoon], Chairman, panel members. My name is Samuel Osei Adu "
    "and this is my MPhil defense on a multi-source context validation framework "
    "for adaptive Zero Trust MFA, supervised by Dr. Kornyo Oliver. Over the next "
    "[X] minutes I'll walk through the problem I set out to solve, how I solved it, "
    "how I evaluated it, and what I found — including being upfront about what the "
    "evaluation did and didn't establish. Let's begin.")
n += 1

# 2. OUTLINE
s = content_slide("Outline", n); n += 1
bullets(s, [
    "Background and Problem Statement",
    "Research Aim, Objectives, and Questions",
    "Related Work and Identified Gaps",
    "Proposed Framework",
    "Experimental Setup",
    "Results",
    "Discussion and Limitations",
    "Contributions and Conclusion",
], size=19, gap=14)
notes(s,
    "Quick roadmap. I'll spend a bit of time on motivation and the gap in the "
    "literature, then move fairly quickly through the framework design since the "
    "detail is in the thesis document, and spend the most time on results — "
    "particularly one finding that explains almost everything else in the results "
    "section, which I'll flag when we get there. I'll close with limitations and "
    "contributions. Happy to take questions at the end, or at any point if the "
    "panel prefers.")

# 3. BACKGROUND
s = content_slide("Background", n, eyebrow="INTRODUCTION"); n += 1
bullets(s, [
    "Remote and hybrid work has expanded the enterprise attack surface — endpoints are often unmanaged, irregularly patched, or on untrusted networks.",
    "Traditional perimeter defences (VPNs, firewalls) permit lateral movement once an attacker authenticates.",
    "Zero Trust Architecture (NIST SP 800-207) requires continuous verification of every access request, regardless of network location.",
    "Adaptive MFA strengthens this by adjusting authentication challenge intensity based on contextual risk signals.",
], size=18, gap=16)
notes(s,
    "The starting point is simple: remote work broke the assumption that "
    "'inside the corporate network' means 'safe.' Once an attacker gets past a "
    "VPN login, there's often nothing stopping lateral movement. Zero Trust is "
    "the industry's answer — verify continuously, trust nothing by default — and "
    "adaptive MFA is how that gets applied at the authentication layer specifically: "
    "instead of one fixed login challenge, the system adjusts based on context. "
    "That context is where my research problem lives — keep this in mind as I "
    "move to the next slide.")

# 4. PROBLEM STATEMENT
s = content_slide("Problem Statement", n, eyebrow="INTRODUCTION"); n += 1
bullets(s, [
    "Contextual signals (GPS, IP, Wi-Fi BSSID, device posture, TLS fingerprint) are consumed by adaptive MFA without validating their reliability.",
    "VPN tunnelling, dynamic IPs, and rogue access points routinely distort these signals in remote environments.",
    "Unvalidated signals inflate risk scores, causing unnecessary step-up challenges for legitimate users.",
    "SIEM systems remain operationally siloed from live MFA enforcement — a temporal gap between detection and response.",
    ("The core problem: no systematic mechanism validates, quality-weights, and integrates heterogeneous contextual signals with real-time security intelligence before authentication enforcement.", 0),
], size=17, gap=14)
notes(s,
    "Here's the gap. Adaptive MFA systems trust their contextual signals at face "
    "value. But in the real world, GPS can be stale, a VPN can make someone look "
    "like they're in another country, a Wi-Fi access point can be spoofed. If the "
    "system doesn't check whether a signal is fresh, internally consistent, or "
    "corroborated by other signals, it either over-trusts a spoofed signal — a "
    "security failure — or over-reacts to a noisy-but-legitimate one — a usability "
    "failure. And separately, SIEM systems that DO detect real threats sit outside "
    "the authentication loop entirely — an analyst might catch something an hour "
    "after the damage is done. That's the problem this thesis addresses on both "
    "fronts.")

# 5. AIM & OBJECTIVES
s = content_slide("Research Aim and Objectives", n, eyebrow="INTRODUCTION"); n += 1
add_textbox(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.9),
            "Aim: design, implement, and evaluate a multi-source context validation "
            "framework that improves authentication accuracy through systematic signal "
            "validation, reducing false positives without compromising security.",
            size=17, italic=True, color=NAVY)
bullets(s, [
    "Design a validation model cross-verifying GPS, IP, Wi-Fi BSSID, device posture, and TLS fingerprint.",
    "Develop a quality-weighted risk integration approach reducing false-positive challenges.",
    "Integrate real-time SIEM intelligence into Zero Trust authentication workflows.",
    "Implement the framework as a modular, containerised architecture.",
    "Empirically evaluate against baseline frameworks on accuracy, usability, latency, and privacy.",
], top=Inches(2.6), size=16, gap=12)
notes(s,
    "The aim, in one sentence: build something that validates context before "
    "trusting it, and prove that doing so actually improves accuracy without "
    "wrecking usability. The five objectives on screen map directly onto the "
    "five chapters of methodology and the five hypotheses I tested — design the "
    "validation model, build the quality-weighted scoring, wire in SIEM, "
    "implement it as real running software rather than a paper design, and then "
    "actually measure it against baselines. I'll go through each of those in "
    "turn.")

# 6. RESEARCH QUESTIONS
s = content_slide("Research Questions", n, eyebrow="INTRODUCTION"); n += 1
bullets(s, [
    "RQ1: How does multi-source validation of contextual signals affect the accuracy of risk-based MFA decisions?",
    "RQ2: To what extent does quality-weighted integration reduce false-positive challenges compared to existing frameworks?",
    "RQ3: How does real-time SIEM-derived intelligence influence adaptive access control under varying threat conditions?",
    "RQ4: What performance overhead does the framework introduce under realistic and constrained network conditions?",
    "RQ5: How does the framework balance security, usability, and privacy?",
], size=17, gap=16)
notes(s,
    "Five research questions, each with a corresponding hypothesis I'll come back "
    "to in the results. Briefly: does validation help accuracy (RQ1), does "
    "quality-weighting specifically cut false positives (RQ2), does SIEM "
    "integration help under active threat (RQ3), what's the latency cost (RQ4), "
    "and does the framework actually balance all three concerns — security, "
    "usability, privacy — rather than trading one off against the others (RQ5). "
    "I'll answer each directly when I get to results, including where the answer "
    "is 'yes, but only partially' — I want to be upfront about that now rather "
    "than have it come up only under questioning.")

# 7. RELATED WORK / GAPS
s = content_slide("Related Work — Identified Gaps", n, eyebrow="LITERATURE REVIEW"); n += 1
bullets(s, [
    "Ahmadi (2025): AI-driven behavioural analytics, Mahalanobis-distance anomaly detection — no multi-source cross-validation, no published thresholds.",
    "Phani Kumar Kanuri (2025): Modular ZTA with Context/Trust Engine — no real-time SIEM integration, no published thresholds.",
    "Jimmy (2025): Context-aware MFA (CAMFA) — publishes no risk-scoring formula; discussed as related work, excluded from quantitative comparison.",
    ("Converging gap: no existing framework combines multi-source signal cross-validation, quality-weighted scoring, real-time SIEM integration, and privacy-aware handling in one pipeline.", 0),
], size=17, gap=14)
notes(s,
    "I selected three of the most closely related published frameworks. Two of "
    "them — Ahmadi and Phani Kumar Kanuri — publish an actual risk-scoring "
    "equation, which meant I could faithfully re-implement them and compare "
    "quantitatively later in this talk. The third, Jimmy's CAMFA framework, is "
    "genuinely relevant conceptually but the paper never publishes a formula to "
    "reproduce — so rather than approximate something and call it a fair "
    "comparison, I've been explicit that it's discussed as related work only "
    "and excluded from the quantitative results. If asked why: I'd rather be "
    "transparent about that limitation than present a shaky re-implementation as "
    "equivalent to the other two.")

# 8. FRAMEWORK ARCHITECTURE
s = content_slide("Proposed Framework — Architecture", n, eyebrow="METHODOLOGY"); n += 1
add_pic(s, os.path.join(FIGURES_DIR, "Figure_3.1_Proposed_Framework_Architecture.png"),
        Inches(2.4), Inches(1.5), height=Inches(5.5))
notes(s,
    "This is the high-level architecture. Four pieces: telemetry collectors "
    "gather the raw signals; the validation layer — the core contribution — "
    "scores each signal's trustworthiness; the risk scoring and policy engine "
    "turns validated signals into a decision; and the authentication gateway "
    "enforces that decision while SIEM feeds back into the loop in real time. "
    "It's built as independent microservices deliberately, so any one piece can "
    "be scaled or temporarily bypassed without taking the whole system down. "
    "I'll go through the two most important pieces — validation and scoring — "
    "on the next two slides.")

# 9. SIGNAL VALIDATION LAYER
s = content_slide("Contextual Signal Validation Layer", n, eyebrow="METHODOLOGY"); n += 1
add_textbox(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.6),
            "Qₛ = Fₛ × Cₛ × Eₛ", size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
bullets(s, [
    "Fₛ — Freshness: decays with signal age.",
    "Cₛ — Consistency: agreement between GPS, IP geolocation, and Wi-Fi BSSID.",
    "Eₛ — Enrichment trust: threat-intelligence penalties (VPN, Tor, malicious IP).",
    ("Signal weights are computed dynamically per session — not a fixed, pre-optimised vector — starting equal and reduced for missing, stale, or mutually inconsistent signals.", 0),
], top=Inches(2.4), size=17, gap=14)
notes(s,
    "This formula is the core contribution. Every signal gets a quality score "
    "that's the product of three things: how fresh it is, how well it agrees "
    "with other independent signals, and how clean it is against threat "
    "intelligence. The consistency term is doing the heaviest lifting for "
    "security specifically — if someone spoofs their GPS but their IP and "
    "Wi-Fi access point still say they're somewhere else, that disagreement "
    "itself becomes the signal. One honest note I want to flag proactively: the "
    "weighting across signal types is dynamic and equal-based, not a "
    "fixed, formally-optimised vector — I chose that heuristically, and I say so "
    "directly in the thesis rather than overclaiming a tuning process I didn't "
    "actually run.")

# 10. RISK SCORING
s = content_slide("Risk Scoring and Policy Engine", n, eyebrow="METHODOLOGY"); n += 1
add_textbox(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.6),
            "R = Rbase + Ranomaly + RSIEM", size=24, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
bullets(s, [
    "Rbase — aggregates per-signal risk weighted by (1 − Qₛ).",
    "Ranomaly — binary anomaly flags from validation (e.g. device/TLS mismatch).",
    "RSIEM — SIEM alert severity (high = 0.30, medium = 0.15).",
    ("Policy thresholds — Allow if R < 0.30; Step-up if 0.30 ≤ R < 0.75; Deny if R ≥ 0.75 — derived from a real ROC sweep (AUC = 0.968), not assumed in advance.", 0),
], top=Inches(2.4), size=17, gap=14)
add_pic(s, os.path.join(FIGURES_DIR, "Figure_3.17_ROC_Analysis_Decision_Thresholds_proposed.png"),
        Inches(8.7), Inches(1.4), height=Inches(5.6))
notes(s,
    "The risk score combines three additive terms — base signal risk, anomaly "
    "flags from validation, and SIEM severity — into one number, and that number "
    "is compared against two thresholds to decide allow, step-up, or deny. Here's "
    "something I want to be upfront about, because it came up during my own "
    "investigation: the original thresholds I started with — 0.25 and 0.75 — were "
    "not actually justified by any real analysis, even though early drafts of "
    "this work claimed they came from ROC analysis. When I actually ran that ROC "
    "sweep for real, on live data, the correct allow threshold turned out to be "
    "0.30, not 0.25 — shown on the right. I'm flagging this proactively because "
    "it's exactly the kind of empirical rigor this thesis is built on: I didn't "
    "just accept a plausible-sounding number, I measured it.")

# 11. DATASETS & BASELINES
s = content_slide("Datasets and Baseline Re-Implementation", n, eyebrow="EXPERIMENTAL SETUP"); n += 1
bullets(s, [
    "CIC-IDS2018 — labelled network attack traffic driving STRIDE-category injection.",
    "RBA (Wiefling et al.) — real production login data with genuine account-takeover ground truth, used for the Spoofing category.",
    "WiGLE (Wi-Fi geolocation) and GeoLite2 (IP geolocation) for cross-validation.",
    ("Ahmadi (2025) and Phani Kumar Kanuri (2025) re-implemented faithfully from their published equations — neither publishes threshold values, so thresholds were calibrated empirically against this study's own data.", 0),
    ("Jimmy (2025) excluded from quantitative comparison — no published risk-scoring formula.", 0),
], size=17, gap=13)
notes(s,
    "For datasets: CIC-IDS2018 is a real, public, widely-cited network-intrusion "
    "dataset — I use it to drive STRIDE-category attack scenarios. I also added "
    "the RBA dataset specifically for the Spoofing category, because it's real "
    "production login data with genuine account-takeover labels, not a simulated "
    "attack — that gives the Spoofing results an extra layer of real-world "
    "grounding. For the baselines: I transcribed Ahmadi's and Phani's equations "
    "directly from their papers and verified the transcription against the "
    "original PDFs. Neither paper publishes numeric thresholds, so — consistent "
    "with what I just said about the ROC sweep — I calibrated those "
    "empirically too, the same way I calibrated my own framework's thresholds, "
    "rather than guessing.")

# 12. EVALUATION PROTOCOL
s = content_slide("Evaluation Protocol", n, eyebrow="EXPERIMENTAL SETUP"); n += 1
bullets(s, [
    "n = 2,054 live sessions per configuration (proposed, ablation, Ahmadi, Phani), identical session stream.",
    "TPR, FPR, Precision, F1, and AUC computed from real classification outcomes.",
    "McNemar's test — the correct test for paired binary classification outcomes on matched sessions — used for statistical significance, not a t-test.",
    "Ablation configuration: proposed framework with the validation layer disabled, isolating its contribution.",
], size=18, gap=16)
notes(s,
    "Every configuration — my framework, the ablation variant, and both "
    "baselines — is evaluated on the exact same 2,054 sessions, submitted in "
    "parallel, so the comparison is apples-to-apples. I want to flag one "
    "specific correction here too: an earlier draft of this evaluation "
    "described a 60/20/20 train-test split with repeated cross-validation "
    "trials — that protocol was never actually implemented. What I actually ran "
    "is a single large-sample live evaluation, which I now describe accurately. "
    "For significance testing, I use McNemar's test rather than a t-test, "
    "because the outcome I'm comparing — correct or incorrect classification — "
    "is paired binary data on matched sessions, and a t-test would be the wrong "
    "tool for that.")

# 13. RESULTS: SECURITY ACCURACY
s = content_slide("Results — Security Accuracy", n, eyebrow="RESULTS"); n += 1
add_table(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(2.2),
    headers=["Metric", "Proposed", "Ablation", "Ahmadi (2025)", "Phani (2025)"],
    rows=[
        ["TPR", "88.30%", "52.13%", "21.86%", "6.11%"],
        ["FPR", "2.86%", "20.95%", "5.71%", "1.90%"],
        ["Precision", "99.83%", "97.88%", "98.61%", "98.35%"],
        ["F1-Score", "0.937", "0.680", "0.358", "0.115"],
        ["AUC", "0.968", "—", "0.572", "0.582"],
    ], font_size=15)
add_textbox(s, Inches(0.7), Inches(4.1), Inches(11.9), Inches(0.5),
            "n = 2,054 sessions per configuration. McNemar's test: p < 0.001 vs every comparison configuration.",
            size=13, italic=True, color=GREY)
add_pic(s, os.path.join(FIGURES_DIR, "Figure_4.1_Security_Accuracy_Metrics.png"),
        Inches(1.2), Inches(4.6), width=Inches(11.0))
notes(s,
    "Here are the headline numbers. My framework: 88.3% detection rate, under 3% "
    "false-positive rate, F1 of 0.937. Both baselines score much lower — Ahmadi "
    "at roughly 22% TPR, Phani at about 6%. I expect the panel's first question "
    "here will be 'why are the baselines so low' — and I have a direct answer on "
    "the next slide, so please hold that question if you can, or I'm happy to "
    "jump ahead. Short version: it's not that my re-implementation is unfair to "
    "them — it's a structural, explainable consequence of what signals their "
    "published equations actually read.")

# 14. RESULTS: PER-STRIDE BREAKDOWN — key slide
s = content_slide("Results — Why the Gap? Detection by STRIDE Category", n, eyebrow="RESULTS — KEY FINDING"); n += 1
add_pic(s, os.path.join(FIGURES_DIR, "Figure_4.5_Detection_Rate_by_STRIDE_Category.png"),
        Inches(1.0), Inches(1.4), width=Inches(11.3))
notes(s,
    "This is the single most important slide in the results section — it "
    "explains almost everything else. I broke detection down by STRIDE attack "
    "category instead of just reporting one aggregate number. Look at Spoofing, "
    "the leftmost group: Ahmadi gets 72%, Phani gets 24% — both reasonably "
    "competent. Now look at every other category — Tampering, Denial of "
    "Service, Elevation of Privilege, Information Disclosure — both baselines "
    "collapse to single digits. That's not random. Ahmadi's and Phani's "
    "published equations only ever read device posture, location, and "
    "time-of-day. Spoofing is the one attack type that actually shows up in "
    "those signals. The other four attack types are network and protocol-layer "
    "attacks — invisible to those equations by construction, no matter how well "
    "I implemented them. My framework detects across all six categories because "
    "it also reads TLS fingerprints and SIEM-correlated network telemetry.")

# 15. RESULTS: INTERPRETATION
s = content_slide("Interpreting the STRIDE Breakdown", n, eyebrow="RESULTS"); n += 1
bullets(s, [
    "Ahmadi and Phani's equations read only device posture, location, and time-of-day signals.",
    "Both detect Spoofing reasonably (72% and 24%) — the one category their signals can observe.",
    "Both perform near-randomly (1-11%) on Tampering, DoS, EoP, and Information Disclosure — network/protocol-layer attacks invisible to their signal scope by construction.",
    ("This is not an implementation weakness — it is a structural, expected consequence of each paper's published equation, verified directly against the source PDFs.", 0),
    ("The proposed framework's multi-source design (TLS, SIEM correlation, geographic cross-validation) detects across all six STRIDE categories.", 0),
], size=17, gap=14)
notes(s,
    "Just to restate the previous slide's point in words, in case the panel "
    "wants it summarised without the chart in front of them: the gap between my "
    "framework and the baselines isn't because my re-implementation of their "
    "equations is unfair or weak — it's an inherent scope limitation of what "
    "signals those equations were designed to read. I verified this by reading "
    "both source papers in full and checking their equations against what I "
    "implemented, line by line. This finding is, if anything, favourable to the "
    "baseline authors — it shows their published work does what it claims within "
    "its stated scope. My contribution is broadening that scope.")

# 16. RESULTS: PERFORMANCE
s = content_slide("Results — Performance", n, eyebrow="RESULTS"); n += 1
stat_tile(s, Inches(0.7), Inches(1.6), Inches(3.6), Inches(1.6), "47 ms", "Median latency", BLUE)
stat_tile(s, Inches(4.55), Inches(1.6), Inches(3.6), Inches(1.6), "2,142 ms", "p95 latency", RED)
stat_tile(s, Inches(8.4), Inches(1.6), Inches(3.6), Inches(1.6), "13-14 ms", "Baselines (median)", GREY)
bullets(s, [
    "Median latency is low; p95 is substantially higher — reflecting external enrichment calls (GeoIP, WiGLE, SIEM) the single-hop baselines never make.",
    "This is variability, not a fixed per-request overhead — an honest characterisation of multi-source validation's real cost.",
], top=Inches(3.6), size=17, gap=14)
add_pic(s, os.path.join(FIGURES_DIR, "Figure_4.2_Performance_Latency_Network_Conditions.png"),
        Inches(0.9), Inches(4.6), width=Inches(11.5))
notes(s,
    "On performance, I want to be direct rather than spin this favourably: "
    "median latency, 47 milliseconds, is genuinely fast. But the 95th "
    "percentile is over two seconds — that's the real cost of calling out to "
    "GeoIP, WiGLE, and SIEM correlation, none of which the baselines do at all "
    "since they only score signals already in the request. An earlier draft of "
    "this thesis claimed a constant 28-millisecond overhead — that number wasn't "
    "real, and once I actually measured the distribution I found it's highly "
    "variable, not constant. I'd rather present the honest, messier finding than "
    "a clean but false one. This is a genuine limitation worth discussing openly "
    "with the panel — I'll return to it in the hypothesis evaluation.")

# 17. RESULTS: STATISTICAL VALIDATION
s = content_slide("Results — Statistical Validation", n, eyebrow="RESULTS"); n += 1
add_table(s, Inches(1.2), Inches(1.7), Inches(10.9), Inches(2.0),
    headers=["Comparison", "Test", "Statistic", "p-value"],
    rows=[
        ["Proposed vs Ablation", "McNemar's (χ²)", "χ² = 557.5", "p < 0.001"],
        ["Proposed vs Ahmadi (2025)", "McNemar's (χ²)", "χ² = 1261.9", "p < 0.001"],
        ["Proposed vs Phani (2025)", "McNemar's (χ²)", "χ² = 1585.1", "p < 0.001"],
    ], font_size=16)
bullets(s, [
    "McNemar's test — paired binary classification outcomes on matched, identical sessions.",
    "Replaces the thesis draft's originally fabricated paired t-test results.",
    "All three comparisons significant at p < 0.001.",
], top=Inches(4.2), size=17, gap=14)
notes(s,
    "To put the accuracy gap on firmer statistical footing: I ran McNemar's "
    "test comparing my framework against each of the three other "
    "configurations, on the identical matched sessions. All three come back "
    "significant at p less than 0.001. I'll mention one more piece of intellectual "
    "honesty here: an earlier version of this chapter reported paired t-test "
    "results with specific p-values that were never actually computed from real "
    "data. I caught that during my own review, removed it, and replaced it with "
    "this real analysis using the statistically correct test for this kind of "
    "paired binary outcome data.")

# 18. RESULTS: ABLATION
s = content_slide("Results — Ablation Analysis", n, eyebrow="RESULTS"); n += 1
add_table(s, Inches(1.7), Inches(1.7), Inches(9.9), Inches(1.6),
    headers=["Configuration", "TPR", "FPR", "F1-Score"],
    rows=[
        ["Full Framework", "88.30%", "2.86%", "0.937"],
        ["Validation Layer Disabled", "52.13%", "20.95%", "0.680"],
    ], font_size=17)
bullets(s, [
    "Disabling validation drops TPR by 36.2 points and raises FPR by 18.1 points.",
    "The validation layer's own contribution exceeds either baseline's entire modelled contribution.",
    ("Granular per-component ablation (isolating TLS, SIEM, or geographic checks individually) was not performed this cycle — noted as future work, not fabricated.", 0),
], top=Inches(3.7), size=17, gap=14)
notes(s,
    "This isolates what the validation layer itself contributes: turn it off, "
    "and TPR drops by 36 points while FPR rises by 18 points — the validation "
    "layer alone accounts for a bigger swing than either baseline's whole "
    "approach. I want to flag a scope limit honestly here too: I only measured "
    "this one on/off configuration. A more granular study — turning off just "
    "the TLS check, or just SIEM, one at a time — would tell us which specific "
    "component matters most, and I didn't run that this cycle. It's explicitly "
    "future work, not something I'm claiming to have measured.")

# 19. DISCUSSION
s = content_slide("Discussion", n, eyebrow="DISCUSSION"); n += 1
bullets(s, [
    "The framework's advantage traces directly to signal scope, not just to its risk formula being \"better\" in the abstract.",
    "Baseline papers' own self-reported figures (92.7%, 96.8%) are on private, unreleased data — not directly comparable to this study's independently verified results.",
    "This study's contribution: the first controlled, statistically validated comparison of these approaches on a shared, disclosed dataset.",
    ("Several claims that would strengthen this thesis further were not established by the experiments performed — reported as open, not resolved.", 0),
], size=17, gap=15)
notes(s,
    "Stepping back: I want to pre-empt a question the panel may well ask — "
    "'Ahmadi's own paper claims 92.7% accuracy, yours shows their "
    "re-implementation at 22%, isn't that a contradiction?' It isn't, and here's "
    "why: their 92.7% is self-reported on their own private simulated dataset "
    "that was never released, so nobody can verify it independently. My 22% is "
    "the same equation, faithfully transcribed, tested against a real, public, "
    "disclosed dataset that anyone can check. Those numbers are answering "
    "different questions and aren't in tension with each other — mine is simply "
    "a stronger standard of evidence.")

# 20. LIMITATIONS
s = content_slide("Limitations", n, eyebrow="DISCUSSION"); n += 1
bullets(s, [
    "CIC-IDS2018's network/protocol-layer attack taxonomy limits detectability for any context-validation framework — a dataset-fit limitation, not a framework weakness.",
    "Endpoint telemetry simulated, not from real devices.",
    "Single-researcher evaluation without independent replication.",
    "Not measured this cycle: granular per-component ablation, SIEM's specific TPR contribution, an independent privacy-leakage audit, and updated network-condition figures.",
], size=17, gap=15)
notes(s,
    "I'd rather present these limitations myself than have them surfaced only "
    "under questioning. The dataset limitation is the most consequential — "
    "CIC-IDS2018 is network-layer-heavy, which caps what any signal-based "
    "authentication framework can detect from certain attack types, mine "
    "included in principle, though my broader signal scope covers more of it "
    "than the baselines'. Endpoint telemetry is simulated. I'm the sole "
    "researcher, so there's no independent replication yet. And several "
    "specific numbers I could have reported — a SIEM-specific accuracy gain, a "
    "formal privacy audit result — I chose not to report, because I didn't "
    "actually run those experiments, even though doing so would have made the "
    "results look stronger on paper.")

# 21. CONTRIBUTIONS
s = content_slide("Research Contributions", n, eyebrow="CONCLUSION"); n += 1
bullets(s, [
    "Theoretical: signal quality (Qₛ = Fₛ × Cₛ × Eₛ) as a first-class variable in adaptive MFA risk computation.",
    "Architectural: direct coupling of SIEM correlation into live authentication enforcement, closing the detection-response gap.",
    "Empirical: a controlled head-to-head comparison with every threshold empirically justified, plus a per-STRIDE-category breakdown explaining precisely where and why the framework outperforms each baseline.",
], size=18, gap=18)
notes(s,
    "Three contributions, one each on the theoretical, architectural, and "
    "empirical fronts. Theoretically, treating signal quality as a first-class "
    "variable — not just present/absent — is, to my knowledge, not something "
    "the existing adaptive-MFA literature does explicitly. Architecturally, "
    "closing the SIEM-to-enforcement gap in real time, not just detecting and "
    "reporting after the fact. And empirically, I'd highlight the "
    "methodological contribution as much as the numbers themselves: every "
    "threshold in this comparison, mine and both baselines', is empirically "
    "derived and disclosed as such, and the per-STRIDE breakdown means the "
    "comparison explains itself rather than just asserting a winner.")

# 22. FUTURE WORK
s = content_slide("Future Work", n, eyebrow="CONCLUSION"); n += 1
bullets(s, [
    "Live enterprise evaluation with real device diversity and network heterogeneity.",
    "Authentication-native ground truth across all STRIDE categories (e.g. LANL or CERT insider-threat datasets).",
    "A genuine session-continuity experiment (multi-request session tracking).",
    "An independent, adversarial privacy-leakage audit.",
    "Granular, component-level ablation study.",
], size=18, gap=16)
notes(s,
    "These five items map directly onto the limitations I just described — "
    "each open question gets a concrete next step rather than being left "
    "vague. If the panel asks 'what would you do with another six months,' "
    "this slide is my answer: live deployment data, richer ground truth for "
    "the categories still reliant on CIC-IDS2018 alone, a real "
    "multi-request session model, an adversarial privacy test, and the "
    "granular ablation I mentioned earlier.")

# 23. CONCLUSION
s = content_slide("Conclusion", n, eyebrow="CONCLUSION"); n += 1
stat_tile(s, Inches(0.6), Inches(1.6), Inches(2.85), Inches(1.5), "88.3%", "TPR", BLUE)
stat_tile(s, Inches(3.6), Inches(1.6), Inches(2.85), Inches(1.5), "2.86%", "FPR", GREEN)
stat_tile(s, Inches(6.6), Inches(1.6), Inches(2.85), Inches(1.5), "0.937", "F1-Score", BLUE)
stat_tile(s, Inches(9.6), Inches(1.6), Inches(2.85), Inches(1.5), "0.968", "AUC", GREEN)
bullets(s, [
    "Validating, quality-weighting, and integrating contextual signals before they influence authentication outcomes produces measurable, statistically significant improvements across every security accuracy metric.",
    "Outperforms both re-implemented published baselines and an ablation configuration under identical, disclosed experimental conditions.",
], top=Inches(3.6), size=18, gap=16)
notes(s,
    "To close: the central claim of this thesis is narrow but well-supported — "
    "validating contextual signals before trusting them, and weighting them by "
    "quality rather than treating them as binary, measurably improves both "
    "detection and false-positive suppression, verified against a real dataset "
    "under a fair, identical, statistically tested comparison against two "
    "published baselines and an ablation control. I've tried throughout this "
    "defense to be as clear about what I didn't establish as what I did — I "
    "think that's what makes the numbers I am presenting trustworthy. Thank "
    "you, and I'm happy to take questions.")

# 24. THANK YOU
s = add_slide(); bg(s, NAVY)
add_textbox(s, Inches(1.0), Inches(3.1), Inches(11.3), Inches(1.2),
            "Thank You", size=44, bold=True, color=WHITE)
add_textbox(s, Inches(1.0), Inches(4.1), Inches(11.3), Inches(0.6),
            "Questions and Discussion", size=20, color=RGBColor(0xC7, 0xD6, 0xEC))
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.85), Inches(2.5), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
notes(s,
    "[Q&A prep, not to be read aloud] Likely questions to anticipate: (1) Why "
    "is p95 latency so high, and is that acceptable for production? — answer "
    "on slide 16, be ready to discuss deployment tradeoffs. (2) Why exclude "
    "Jimmy entirely rather than approximate it? — slide 7, no formula to "
    "reproduce faithfully. (3) Isn't a single-researcher, simulated-telemetry "
    "evaluation a weak form of validation? — slide 20, agree directly, point "
    "to future work slide 22. (4) Why McNemar's and not a t-test? — slide 12/17, "
    "paired binary outcome data. (5) How were baseline thresholds chosen if "
    "the papers don't publish them? — slide 11, same empirical ROC-sweep "
    "methodology used for my own thresholds, disclosed as such.")

prs.save(OUT_PATH)
print(f"Saved {n} content slides + title/thank-you to: {OUT_PATH}")
